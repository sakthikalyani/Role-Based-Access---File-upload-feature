import os
import pypdf
import docx
from pptx import Presentation
import pandas as pd
from pathlib import Path
import markdown
from bs4 import BeautifulSoup

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
    return text

def extract_text_from_docx(file_path):
    """Extract text from Word document"""
    text = ""
    try:
        doc = docx.Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
    return text

def extract_text_from_csv(file_path):
    """Convert CSV to markdown format"""
    text = ""
    try:
        df = pd.read_csv(file_path)
        text = df.to_markdown(index=False)
    except Exception as e:
        print(f"Error reading CSV {file_path}: {e}")
    return text

def extract_text_from_markdown(file_path):
    """Extract text from markdown file"""
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            md_content = file.read()
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
    except Exception as e:
        print(f"Error reading Markdown {file_path}: {e}")
    return text

def process_document(file_path):
    """
    Process any supported document type and return text content
    Supported formats: PDF, DOCX, PPTX, CSV, MD
    """
    file_extension = Path(file_path).suffix.lower()

    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif file_extension in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif file_extension == '.csv':
        return None
    elif file_extension in ['.xlsx', '.xls']:
        # same for excel
        return None
    elif file_extension == '.md':
        return extract_text_from_markdown(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

def save_as_markdown(text_content, output_path, metadata=None):
    """Save extracted text as markdown with optional metadata"""
    with open(output_path, 'w', encoding='utf-8') as f:
        if metadata:
            f.write("---\n")
            for key, value in metadata.items():
                f.write(f"{key}: {value}\n")
            f.write("---\n\n")
        f.write(text_content)
